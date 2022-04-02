import sys
import os
import comtypes.client


import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import os, sys
import matplotlib
import seaborn as sns
from matplotlib import cm
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter
from matplotlib.patches import Circle, Wedge, Rectangle

import plotly
import plotly.graph_objects as go

import json

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Pt
from pptx.util import Inches

from babel.numbers import format_currency

mypres = Presentation("Investor Wealth Report-v3.pptx")

# met_img_path = 'metMadhur Jindal9188006668882022-03-29145955'+'.png'
# bar_img_path = 'barMadhur Jindal9188006668882022-03-29145955'+'.png'
# lin_img_path = 'lineMadhur Jindal9188006668882022-03-29145955'+'.png'

# slide0 = mypres.slides[0] 
# slide1 = mypres.slides[1]
# slide2 = mypres.slides[2]

# bar_left = Inches(0.1)
# bar_top = Inches(8.5) 
    
# bar_height = Inches(0.6) 
    
# bar_pic = slide0.shapes.add_picture(bar_img_path, bar_left,
#                                 bar_top, height = bar_height)
# # For margins
# met_left = Inches(0.1)
# met_top = Inches(2.835) 
    
# met_height = Inches(1.45) 
    
# met_pic = slide0.shapes.add_picture(met_img_path, met_left,
#                                 met_top, height = met_height)

# lin_left = Inches(0.2)
# lin_top = Inches(7.2) 

# lin_height = Inches(3.0) 
    
# line_pic = slide1.shapes.add_picture(lin_img_path, lin_left,
#                                 lin_top, height = lin_height)



# name = "dinesh"
# mobile = "99999999"
# net_worth_as_input = "INR 20-25 laks"
# age = "28"
# life_stage = "Single"
# savings_as_input = "20%"
# income_source = "Stable"
# income_as_input = "INR 10-15 lakh"

# def fill_bio(mypres,sno,shapeNo,r,c, text):
#         '''
#         Fills the table in pptx

#         '''
#         p = mypres.slides[sno].shapes[shapeNo].table.cell(r,c).text_frame.paragraphs[0]
#         #clears the cell
#         run = p.clear()
#         run = p.add_run()
#         run.text = text
#         font = run.font
#         font.name = 'Times New Roman'
#         font.size = Pt(8)

# fill_bio(mypres,0,20,0,1, name)
# fill_bio(mypres,0,20,1,1, '+'+str(mobile)[:2]+str('-'+mobile[2:]))
# fill_bio(mypres,0,20,2,1, net_worth_as_input)                
# fill_bio(mypres,0,20,0,3, str(age)+' years')
# fill_bio(mypres,0,20,1,3, life_stage)
# fill_bio(mypres,0,20,2,3, savings_as_input)                                   
# fill_bio(mypres,0,20,0,5, income_source)
# fill_bio(mypres,0,20,1,5, income_as_input) 

# fill_bio(mypres,1,14,1,1, str(0.001*100)+"%")
# fill_bio(mypres,1,14,1,2, "4")
# fill_bio(mypres,1,14,2,1, "4")
# fill_bio(mypres,1,14,2,2, "4")
# fill_bio(mypres,1,14,3,1, "4")
# fill_bio(mypres,1,14,3,2, "4")
# fill_bio(mypres,1,14,4,1, "4")
# fill_bio(mypres,1,14,4,2, "4")

# mypres.save('TEST1.pptx')

# p = mypres.slides[0].shapes[21].table.cell(1,1).text_frame.paragraphs[0]
p = mypres.slides[2].shapes[17]

for x in mypres.slides[2].shapes:
    print(x.name,x.shape_id)

print(round(33.33,2))