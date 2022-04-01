import sys
import os
import comtypes.client


import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import os, sys
import re
import matplotlib
import seaborn as sns
from matplotlib import cm
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter
from matplotlib.patches import Circle, Wedge, Rectangle

import plotly
import plotly.graph_objects as go

import json

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Pt
from pptx.util import Inches

from babel.numbers import format_currency



def nest():
    '''
    A nested dictionary containing chatbot's users' possible
    responses along with associated values within the list.
    The first dictionary pairs the Question number or the 
    column number of the Spreadsheet in the database, whereas
    the enclosed dictionary contains the list of associated 
    values with each option
    

    Returns
    -------
    dict
        Nested Dictionary of Questionnaire and values of 
            the user's responses

    '''
    
    
    return {2:{'Not at all': [17.5, 'conservative', 'tolerance'],
    'Maybe': [39.5, 'moderate', 'tolerance'],
    'After considering safety': [49.5, 'balanced', 'tolerance'],
    'Blindly go': [82.5, 'aggressive', 'tolerance']},

    3:{'1-2 years': [17.5, 'conservative', 'tolerance'],
    '3-5 years': [39.5, 'moderate', 'tolerance'],
    '6-10 years': [49.5, 'balanced', 'tolerance'],
    '11-15 years': [59.5, 'assertive (growth)', 'tolerance'],
    'More than 15 years': [82.5, 'aggressive', 'tolerance']},


    4:{'Very inexperienced': [17.5, 'conservative', 'tolerance'],
    'Somewhat inexperienced': [39.5, 'moderate', 'tolerance'],
    'Somewhat experienced': [49.5, 'balanced', 'tolerance'],
    'Experienced': [59.5, 'assertive (growth)', 'tolerance'],
    'Very experienced': [82.5, 'aggressive', 'tolerance']},


    5:{'Sell 100%': [17.5, 'conservative', 'tolerance'],
    'Sell >50%': [39.5, 'moderate', 'tolerance'],
    'Sell <50%': [49.5, 'balanced', 'tolerance'],
    'Sell 0%': [59.5, 'assertive (growth)', 'tolerance'],
    'Buy': [82.5, 'aggressive', 'tolerance']},

    6:{'A real risk avoider': [17.5, 'conservative', 'tolerance'],
    'Cautious': [39.5, 'moderate', 'tolerance'],
    'Willing to take risk': [49.5, 'balanced', 'tolerance'],
    'A Real risk taker': [82.5, 'aggressive', 'tolerance']},

    7:{'I strongly disagree': [82.5, 'aggressive', 'tolerance'],
    'I disagree': [59.5, 'assertive (growth)', 'tolerance'],
    'I somewhat agree': [49.5, 'balanced', 'tolerance'],
    'I agree': [39.5, 'moderate', 'tolerance'],
    'I strongly agree': [17.5, 'conservative', 'tolerance']},

    8:{'Sell 100%': [17.5, 'conservative', 'tolerance'],
    'Sell >50%': [39.5, 'moderate', 'tolerance'],
    'Sell <50%': [49.5, 'balanced', 'tolerance'],
    'Sell 0%': [59.5, 'assertive (growth)', 'tolerance'],
    'Buy': [82.5, 'aggressive', 'tolerance']},

    9:{'Win Rs. 110': [17.5, 'conservative', 'tolerance'],
    'Win Rs. 120': [39.5, 'moderate', 'tolerance'],
    'Win Rs. 150': [49.5, 'balanced', 'tolerance'],
    'Win Rs. 200': [59.5, 'assertive (growth)', 'tolerance'],
    'Win Rs. 300': [82.5, 'aggressive', 'tolerance']},

    10:{'Sell 100%' : [17.5, 'conservative', 'tolerance'],
    'Sell >50%': [39.5, 'moderate', 'tolerance'],
    'Sell <50%': [49.5, 'balanced', 'tolerance'],
    'Sell 0%' : [59.5, 'assertive (growth)', 'tolerance'],
    'Buy': [82.5, 'aggressive', 'tolerance']},

    11:{'Not comfortable': [17.5, 'conservative', 'tolerance'],
    'A little hesitant': [39.5, 'moderate', 'tolerance'],
    'Reasonably comfortable': [49.5, 'balanced', 'tolerance'],
    'Very comfortable': [82.5, 'aggressive', 'tolerance']},

    12:{'< INR 5 lakh': [17.5, 'conservative', 'capacity'],
    'INR 5-10 lakh': [39.5, 'moderate', 'capacity'],
    'INR 11-20 lakh': [49.5, 'balanced', 'capacity'],
    'INR 21-30 lakh': [59.5, 'assertive (growth)', 'capacity'],
    'INR >30 lakh': [82.5, 'aggressive', 'capacity']},

    13:{'Very unstable': [17.5, 'conservative', 'capacity'],
    'Unstable': [39.5, 'moderate', 'capacity'],
    'Somewhat stable': [49.5, 'balanced', 'capacity'],
    'Stable': [59.5, 'assertive (growth)', 'capacity'],
    'Very stable': [82.5, 'aggressive', 'capacity']},

    14:{'less than 10%': [17.5, 'conservative', 'capacity'],
    '11-25%': [39.5, 'moderate', 'capacity'],
    '26-40%': [49.5, 'balanced', 'capacity'],
    '41-55%': [59.5, 'assertive (growth)', 'capacity'],
    '> 55%': [82.5, 'aggressive', 'capacity']},

    15:{'< INR 10 lakh': [17.5, 'conservative', 'capacity'],
    'INR 10-25 lakh': [39.5, 'moderate', 'capacity'],
    'INR 25-50 lakh': [49.5, 'balanced', 'capacity'],
    'INR 50-100 lakh': [59.5, 'assertive (growth)', 'capacity'],
    '> INR 100 lakh': [82.5, 'aggressive', 'capacity']},

    16:{'Single, No dependent': [82.5, 'aggressive', 'capacity'],
    'Single': [59.5, 'assertive (growth)', 'capacity'],
    'Young Family': [49.5, 'balanced', 'capacity'],
    'Nearing Retirement': [39.5, 'moderate', 'capacity'],
    'Retired': [17.5, 'conservative', 'capacity']}
    }





def income_and_saving(data_row):
    '''
    This function uses user inputs
    about income and savings from 
    data_row 12 and 14, and converts 
    them into numerical figures for
    further usage

    Parameters
    ----------
    data_row : the user's row.

    Returns
    -------
    numerically doable values of income and saving.

    '''
    income = 0 
    saving = 0 
    i,s = data_row[12], data_row[14]
    if i == '< INR 5 lakh':
        income = 250000
    elif i == 'INR 5lac-10 lakh':
        income = 750000
    elif i == 'INR 11-20 lakh':
        income = 1550000
    elif i == 'INR 21-30 lakh':
        income = 2550000
    elif i == 'INR >30 lakh':
        income = 4000000
    if s == 'less than 10%':
        saving = 0.05
    elif s == '11-25%':
        saving = 0.18
    elif s == '26-40%':
        saving = 0.33
    elif s == '41-55%':
        saving = 0.48
    elif s == '> 55%':
        saving = 0.60
    return income, float(saving)





#data_dict = nest()

#print(income_and_saving(data_row))
def scores(data_row, data_dict):
    '''
    Evaluates scores based on the user's inputs

    Parameters
    ----------
    data_row : list
        user's responses.
    data_dict : nested dict
        the data dictionary we get from using nest().

    Returns
    -------
    sum_capacities : int
        sum of all questions described as of explaining the capacity of the user.
    sum_tolerances : int
        sum of all questions described as of explaining the tolerance of the user.
    avg_capacities : int
        average of all questions described as of explaining the capacity of the user.
    avg_tolerances : int
        average of all questions described as of explaining the tolerance of the user.
    total : int
        the total score of the user.
    avg_total : int
        the average score of the user.

    '''
    tolerances, capacities = [], []
    for element in data_row: # data row
        ix = data_row.index(element) # grab the index (i-e col)
        for key in data_dict.keys(): # data dict 1st level
            if ix == key: 
                if data_dict[key][element][2] == 'tolerance':
                    tolerances.append(data_dict[key][element][0])
                elif data_dict[key][element][2] == 'capacity':
                    capacities.append(data_dict[key][element][0])
    sum_capacities, sum_tolerances = sum(tolerances), sum(capacities)
    avg_capacities, avg_tolerances = sum_capacities/5, sum_tolerances/10
    total = sum_capacities + sum_tolerances
    avg_total = total/16
    return sum_capacities, sum_tolerances, avg_capacities, avg_tolerances, total, avg_total


# print(scores(data_row, data_dict))


# sum_capacities, sum_tolerances, avg_capacities, avg_tolerances, total, avg_total = scores(data_row, data_dict)





def risk_profile(avg_total):
    '''
    Defines the ratio between debt and equity
    as debt + equity = 1.00

    Parameters
    ----------
    avg_total : int
        the average score of the user, obtained from `scores()` .

    Returns
    -------
    debt : float
        ratio of debt .
    equity : TYPE
        ratio of equity.

    '''
    if 0 < avg_total <= 35:
        equity = float(0.15)
    elif 35 < avg_total <= 45:
        equity = float(0.30)
    elif 45 < avg_total <= 55:
        equity = float(0.50)
    elif 55 < avg_total <= 65:
        equity = float(0.70)
    else:
        equity = float(0.85)
    debt = round(1 - equity, 3)
    return debt, equity




# print(risk_profile(avg_total))
#debt, equity = risk_profile(avg_total)
#nif, figi, aror = 0.15, 0.07, 0.10
 
def schedule(sip, debt, equity, nif, figi):
    '''
    

    Parameters
    ----------
    sip : int
        SIP amount.
    debt : float
        Debt ratio.
    equity : float
        Equity ratio.
    nif : float
        Nifty Index Fund.
    figi : TYPE
        FI Govt Long Index Fund.

    Returns
    -------
    invested_amount : int
        Amount for Investment funds.
    portfolio_amount : int
        Amount for Portfolio funds.
    years : years
        years, 25 by default.

    '''
    years = [i for i in range(1,26)]
    equity_a, debt_a = [], []
    e = (sip * equity)+((sip*equity/2)*(1+nif))-(sip*equity/2)
    d = (sip*debt)+((sip*debt/2)*(1+figi))-(sip*debt/2)
    equity_a.append(e)
    debt_a.append(d)
    stopper = 0
    while stopper < len(years)-1:
        equity_a.append(e + equity_a[-1]*(1+nif))
        debt_a.append(d + debt_a[-1]*(1+figi))
        stopper+=1
    invested_amount = [sip*year for year in years]
    portfolio_amount = [sum(pair) for pair in zip(equity_a, debt_a)]
    return invested_amount, portfolio_amount, years
#%
#invested_amount, portfolio_amount, years = schedule(sip, debt, equity, nif, figi)

def forecast(invested_amount, portfolio_amount, years, sip, debt, equity, nif, figi):
    '''
    

    Parameters
    ----------
    invested_amount : int
        The amount for investment funds.
    portfolio_amount : int
        The amount for portfolio funds.
    years : int
        years, span.
    sip : int
        SIP income.
    debt : float
        Debt.
    equity : float
        Equity.
    nif : float
        Nifty Index.
    figi : float
        FI Govt Long Index.

    Returns
    -------
    df_forecast : DataFrame 
        Forecast of 25 years.

    '''
    inv, plf, yrs = schedule(sip=sip, debt=debt, equity=equity, nif=nif, figi=figi)

    df_forecast = pd.DataFrame(
        {'year': yrs,
         'invested_amount': inv,
         'portfolio_amount': plf
        })

    df_forecast.iloc[:,1:] = df_forecast.iloc[:,1:]/10**5
    return df_forecast

# print(forecast(sip, invested_amount, portfolio_amount, years))



def tolerance_remarks(score_c, score_t):
    '''
    Generates recommedations depending on the score evaluated by the user's inputs

    Parameters
    ----------
    score_c : int
        Total score of the questions about capacity.
    score_t : int
        Total score of the questions about tolerance.

    Returns
    -------
    String
        Remarks (Recommendation).

    '''
    far_less = 'Your risk tolerance does not match with your risk capacity. This means that the amount of risk you are willing to take is far less than the amount of risk you can afford to take.'
    very_less = 'Your risk tolerance does not match with your risk capacity. This means that the amount of risk you are willing to take is very less as compared to the amount of risk you can afford to take.'
    less = 'Your risk tolerance does not match with your risk capacity. This means that the amount of risk you are willing to take is less than the amount of risk you can afford to take.'
    okayish = 'Your risk tolerance matches with your risk capacity. This means that the amount of risk you are willing to take matches with the amount of risk you need.'
    more = 'Your risk tolerance does not match with your risk capacity. This means that the amount of risk you are willing to take is more than the amount of risk you can afford to take.'
    quite_more = 'Your risk tolerance does not match with your risk capacity. This means that the amount of risk you are willing to take is pretty much more than the amount of risk you can afford to take.' 
    far_more = 'Your risk tolerance does not match with your risk capacity. This means that the amount of risk you are willing to take is far more than the amount of risk you can afford to take.'

    grades = []
    for score in [score_c, score_t]:

        if 0 < score <= 35:
            grade = 'Conservative'
        elif 35 < score <= 45:
            grade = 'Moderate'
        elif 45 < score <= 55:
            grade = 'Balanced'
        elif 55 < score <= 65:
            grade = 'Assertive (Growth)'
        else:
            grade = 'Aggressive'
        grades.append(grade)
    if grades in [['Conservative', 'Balanced'], ['Moderate', 'Assertive (Growth)'], ['Balanced', 'Aggressive']]:
        return more
    elif grades in [['Balanced', 'Conservative'], ['Assertive (Growth)', 'Moderate'], ['Aggressive','Balanced']]:
        return less
    elif grades in [['Conservate', 'Assertive (Growth)'], ['Moderate', 'Aggressive']]:
        return quite_more
    elif grades in [['Assertive (Growth)', 'Conservative'], ['Aggressive', 'Moderate']]:
        return very_less
    elif grades in [['Conservative', 'Aggressive']]:
        return far_more
    elif grades in [['Aggressive', 'Conservative']]:
        return far_less
    else:
        return okayish

def counter(i=[10012022000000]):

    i[0]+=1 
    srn = i[0]
    return srn

# VISUALIZATIONS




def stackbar(debt, equity, name, mobile, code):
    '''
    

    Parameters
    ----------
    equity : float
        equity.
    debt : float
        debt.
    name : str
        name.
    mobile : str
        mobile.
    code : str
        code.

    Returns
    -------
    A Stackbar visualizing debt vs equity is stored in the directory.

    '''

    top_labels = ['Debt',
                  'Equity']

    colors = ['rgba(255,218,223,255)',
              'rgba(191,127,127,255)']

    x_data = [[debt*100,equity*100]]

    y_data = ['Debt/Equity']
    fig = go.Figure()

    for i in range(0, len(x_data[0])):
        for xd, yd in zip(x_data, y_data):
            fig.add_trace(go.Bar(
                x=[xd[i]], y=[yd],
                orientation='h',
                marker=dict(
                    color=colors[i],
                    line=dict(color='rgb(248, 248, 249)', width=5)
                )
            ))

    fig.update_layout(
        xaxis=dict(
            showgrid=False,
            showline=False,
            showticklabels=False,
            zeroline=False,
            domain=[0.15, 1]
        ),
        yaxis=dict(
            showgrid=False,
            showline=False,
            showticklabels=False,
            zeroline=False,
        ),
        barmode='stack',
        paper_bgcolor='rgb(258, 258, 255)',
        plot_bgcolor='rgba(255,255,255,255)',
        margin=dict(l=0, r=0, t=50, b=0),
        width = 1300,
        height = 200,
        showlegend=False,
    )

    annotations = []

    for yd, xd in zip(y_data, x_data):
        annotations.append(dict(xref='x', yref='y',
                                x=xd[0] / 2, y=yd,
                                text=str(xd[0]) + '%',
                                font=dict(family='Arial', size=40,
                                          color='rgb(67, 67, 67)'),
                                showarrow=False))
        if yd == y_data[-1]:
            annotations.append(dict(xref='x', yref='paper',
                                    x=xd[0] / 2, y=1.2,
                                    text=top_labels[0],
                                    font=dict(family='Arial', size=45,
                                              color='rgb(67, 67, 67)'),
                                    showarrow=False))
        space = xd[0]
        for i in range(1, len(xd)):
                annotations.append(dict(xref='x', yref='y',
                                        x=space + (xd[i]/2), y=yd,
                                        text=str(xd[i]) + '%',
                                        font=dict(family='Arial', size=40,
                                                  color='rgb(67, 67, 67)'),
                                        showarrow=False))
                if yd == y_data[-1]:
                    annotations.append(dict(xref='x', yref='paper',
                                            x=space + (xd[i]/2), y=1.2,
                                            text=top_labels[i],
                                            font=dict(family='Arial', size=45,
                                                      color='rgb(67, 67, 67)'),
                                            showarrow=False))
                space += xd[i]

    fig.update_layout(annotations=annotations)


    fig.show()

    fig.write_image('bar'+str(name)+str(mobile)+str(code)+'.png', engine='kaleido')








def degree_range(n): 
    '''
    For Gauge, run Gauge function

    '''
    start = np.linspace(0,180,n+1, endpoint=True)[0:-1]
    end = np.linspace(0,180,n+1, endpoint=True)[1::]
    mid_points = start + ((end-start)/2.)
    return np.c_[start, end], mid_points

def rot_text(ang):
    '''
    For Gauge, run Gauge function

    '''
    rotation = np.degrees(np.radians(ang) * np.pi / np.pi - np.radians(90))
    return rotation

def gauge(name, 
          mobile,
          code,
          labels=['Conservative','Moderate','Balanced','Assertive','Aggressive'],
          colors='jet_r', 
          arrow=1, 
          title='',
          size = (5,3)): 
    '''
    

    Parameters
    ----------
    name : str
        name.
    mobile : str
        mobile.
    code : str
        code.
    labels : list
        The default is ['Conservative','Moderate','Balanced','Assertive','Aggressive'].
    colors : list
        Optional.
    arrow : num
        The needle.
    title : str
        Title or remarks just below the chart.
    size : tuple
        size.


    Returns
    -------
    A Gauge/Speedometer-like chart telling about users category. It is stored in the directory.

    '''
          
    

    
    N = len(labels)
    
    if arrow > N: 
        raise Exception("\n\nThe category ({}) is greated than         the length\nof the labels ({})".format(arrow, N))
 
    

    if isinstance(colors, str):
        cmap = cm.get_cmap(colors, N)
        cmap = cmap(np.arange(N))
        colors = cmap[::-1,:].tolist()
    if isinstance(colors, list): 
        if len(colors) == N:
            colors = colors[::-1]
        else: 
            raise Exception("\n\nnumber of colors {} not equal             to number of categories{}\n".format(len(colors), N))

    
    fig, ax = plt.subplots(figsize=size, dpi=100) # figsize=(5,5)

    ang_range, mid_points = degree_range(N)

    labels = labels[::-1]
    

    patches = []
    for ang, c in zip(ang_range, colors): 
        # sectors
        patches.append(Wedge((0.,0.), .4, *ang, facecolor='w', lw=2))
        # arcs
        patches.append(Wedge((0.,0.), .4, *ang, width=0.10, facecolor=c, lw=2, alpha=0.5))
    
    [ax.add_patch(p) for p in patches]

    
    """
    set the labels (e.g. 'LOW','MEDIUM',...)
    """

    for mid, lab in zip(mid_points, labels): 

        ax.text(0.35 * np.cos(np.radians(mid)), 
                0.35 * np.sin(np.radians(mid)), 
                lab,
                horizontalalignment='center', 
                verticalalignment='center', 
                fontsize=8,
                fontweight='bold', 
                rotation = rot_text(mid))


    r = Rectangle((-0.4,-0.1),0.8,0.1, facecolor='w', lw=2)
    ax.add_patch(r)
    
    ax.text(0, 
            -0.05, 
            title, 
            horizontalalignment='center',
            verticalalignment='center', 
            fontsize=20, 
            fontweight='bold')


    pos = mid_points[abs(arrow - N)]
    
    ax.arrow(0, 
             0, 
             0.225 * np.cos(np.radians(pos)), 
             0.225 * np.sin(np.radians(pos)),
             width=0.02, 
             head_width=0.03, 
             head_length=0.12, 
             fc='w', 
             ec='k')
    
    ax.add_patch(Circle((0, 0), radius=0.02, facecolor='k'))
    ax.add_patch(Circle((0, 0), radius=0.01, facecolor='w', zorder=11))


    
    ax.set_frame_on(False)
    ax.axes.set_xticks([])
    ax.axes.set_yticks([])
    ax.axis('equal')
    plt.tight_layout()
    plt.savefig('met'+str(name)+str(mobile)+str(code)+'.png', dpi=200)

def pointer(avg_total):
    if 0 < avg_total <= 35:
        return 1
    elif 35 < avg_total <= 45:
        return 2
    elif 45 < avg_total <= 55:
        return 3
    elif 55 < avg_total <= 65:
        return 4
    else:
        return 5


#gauge(labels=['Conservative','Moderate','Balanced','Assertive','Aggressive'], colors=["#FFB6C1","#EE6363","#CD5555","#8B3A3A","#800000"], arrow=pointer(avg_total), size=(5,3), title=str('Your Risk Score is {}'.format(int(avg_total))))
    
    
    





def line_chart(df_forecast, name, mobile, code):
    '''
    Generates Line chart of the forecast dataframe

    Parameters
    ----------
    df_forecast : DataFrame ''forecast''
        The forecast dataframe, generated through forecast().
    name : str
        name.
    mobile : str
        mobile.
    code : str
        code for uniqueness.

    Returns
    -------
    Saves a line chart in the directory for the forecast.

    '''
    sns.set_style('ticks')
    fig, ax = plt.subplots(figsize=(15, 7))
    # Create a scatter plot
    sns.lineplot(x='year',
                    y='portfolio_amount',
                    ax=ax,
                    data=df_forecast,
                    palette='rocket', color='#B22222', linewidth=3.0, label='Portfolio Amount')
    sns.lineplot(x='year',
                    y='invested_amount',
                    ax=ax,
                    data=df_forecast,
                    palette='rocket', color='#483D8B', linewidth=3.0, label='Investment Amount')
    right_side = ax.spines["right"]
    right_side.set_visible(False)
    top_side = ax.spines["top"]
    top_side.set_visible(False)
    ax.annotate(text=str(round(df_forecast.iloc[4,2],1))+' Lakhs', xy=(4, df_forecast.iloc[4,2]*1.50), fontsize=12, weight='bold', bbox=dict(boxstyle='circle', facecolor='#B22222',alpha=0.1))
    ax.annotate(text=str(round(df_forecast.iloc[9,2],1))+' Lakhs', xy=(9, df_forecast.iloc[9,2]*1.35), fontsize=12, weight='bold',bbox=dict(boxstyle='circle', facecolor='#B22222', alpha=0.1))
    ax.annotate(text=str(round(df_forecast.iloc[14,2],1))+' Lakhs', xy=(14, df_forecast.iloc[14,2]*1.25), fontsize=12, weight='bold', bbox=dict(boxstyle='circle',facecolor='#B22222', alpha=0.1))
    ax.annotate(text=str(round(df_forecast.iloc[19,2],1))+' Lakhs', xy=(19, df_forecast.iloc[19,2]*1.15), fontsize=12, weight='bold', bbox=dict(boxstyle='circle', facecolor='#B22222', alpha=0.1))
    ax.annotate(text=str(round(df_forecast.iloc[24,2],1))+' Lakhs', xy=(24, df_forecast.iloc[24,2]*1.05), fontsize=12, weight='bold', bbox=dict(boxstyle='circle', facecolor='#B22222', alpha=0.1))

    ax.annotate(text=str(round(df_forecast.iloc[24,1],1))+' Lakhs', xy=(24, df_forecast.iloc[24,1]*1.50), fontsize=12, weight='bold', bbox=dict(boxstyle='circle',facecolor='#483D8B', alpha=0.1))

    our_yticks = list()
    
    if df_forecast.iloc[-1,2] <= 100:
        our_yticks = [j for j in range(int(df_forecast.iloc[-1,2]) + 30) if j%10==0]
    elif 100 < df_forecast.iloc[-1,2] <= 500:
        our_yticks = [j for j in range(int(df_forecast.iloc[-1,2]) + 50) if j%50==0] #largest ytick in range
    elif 500 < df_forecast.iloc[-1,2] <= 1000:
        our_yticks = [j for j in range(int(df_forecast.iloc[-1,2]) + 100) if j%100==0] #largest ytick in range
    elif 1000 < df_forecast.iloc[-1,2] < 5000:
        our_yticks = [j for j in range(int(df_forecast.iloc[-1,2]) + 500) if j%500==0]
    else:
        our_yticks = [j for j in range(int(df_forecast.iloc[-1,2]) + 1000) if j%1000==0]
    
    our_yticklabels = [(str(tick) + ' Lakhs') for tick in our_yticks if str(tick)]

    ax.set_yticks(our_yticks)
    ax.set_yticklabels(our_yticklabels, fontsize=15, fontname='Times New Roman', weight='bold')

    ax.set_xticks([0,5,10,15,20,25])
    ax.set_xticklabels(["Present", "5 years", "10 years", "15 years","20 years", "25 years"], fontsize=15, fontname='Times New Roman', weight='bold')
    plt.legend(loc='upper center',
                  ncol=2, fancybox=True, shadow=True,frameon=False,fontsize=15)
    ax.xaxis.label.set_visible(False)
    ax.yaxis.label.set_visible(False)
    plt.grid(axis='y', lw=0.5, ls=':')




    plt.savefig('line'+str(name)+str(mobile)+str(code)+'.png', dpi=200)    
    
#line_chart(df_forecast)




# TOWARDS POWERPOINT


def pptx_work(presentation, name, mobile, code,date, sip, aror, figi, avg_capacities, avg_tolerances,
              net_worth_as_input, age, life_stage, savings_as_input, income_source, income_as_input, data_row):
    '''
    It does the whole work of modifying/creating the pptx file with requests
    made above. It takes in a lot of params. They are self explanatory and
    mostly derived from the data_row

    Parameters
    ----------
    various but simple

    Returns
    -------
    A pptx file named 'base.pptx' is created/modified in the directory.

    '''
    mypres = Presentation(presentation)
    met_img_path = 'met'+str(name)+str(mobile)+str(code)+'.png'
    bar_img_path = 'bar'+str(name)+str(mobile)+str(code)+'.png'
    lin_img_path = 'line'+str(name)+str(mobile)+str(code)+'.png'

    def fill_bio(mypres,r,c, text):
        '''
        Fills the table in pptx

        '''
        p = mypres.slides[0].shapes[21].table.cell(r,c).text_frame.paragraphs[0]
        #clears the cell
        run = p.clear()
        run = p.add_run()
        run.text = text
        font = run.font
        font.name = 'Times New Roman'
        font.size = Pt(8)

    fill_bio(mypres,0,1, name)
    fill_bio(mypres,1,1, '+'+str(mobile)[:2]+str('-'+mobile[2:]))
    fill_bio(mypres,2,1, net_worth_as_input)                 #(str('INR ')+f"{int(net_worth):,}"))
    fill_bio(mypres,0,3, str(age)+' years')
    fill_bio(mypres,1,3, life_stage)
    fill_bio(mypres,2,3, savings_as_input)                                    #str(savings*100)+'%')
    fill_bio(mypres,0,5, income_source)
    fill_bio(mypres,1,5, income_as_input)                                    #str('INR ')+f"{income:,}")




    def pptx_job(mypres, name, mobile, code, date, avg_capacities, avg_tolerances, sip, figi, aror):

        '''
        Integrates the images
        '''
        
        slide0 = mypres.slides[0] 
        slide1 = mypres.slides[1]
        
        bar_left = Inches(0.1)
        bar_top = Inches(8.5) 
          
        bar_height = Inches(0.6) 
          
        bar_pic = slide0.shapes.add_picture(bar_img_path, bar_left,
                                       bar_top, height = bar_height)
        # For margins
        met_left = Inches(0.1)
        met_top = Inches(2.835) 
          
        met_height = Inches(1.45) 
          
        met_pic = slide0.shapes.add_picture(met_img_path, met_left,
                                       met_top, height = met_height)
    
        lin_left = Inches(0.2)
        lin_top = Inches(2.5) 
    
        lin_height = Inches(3.0) 
          
        line_pic = slide1.shapes.add_picture(lin_img_path, lin_left,
                                       lin_top, height = lin_height)
    
    
    pptx_job(mypres, name, mobile, code, date, avg_capacities, avg_tolerances, sip, figi, aror)


    mypres.slides[1].shapes[12].text_frame.paragraphs[0].text = tolerance_remarks(avg_capacities, avg_tolerances)
    for paragraph in mypres.slides[1].shapes[12].text_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.name = 'Times New Roman'
        
    mypres.slides[1].shapes[10].text_frame.paragraphs[0].text = 'A monthly SIP amount of '+str(format_currency(sip/12, 'INR', locale='en_IN'))+' is recommended based on your current lifestyle and commitments. However, the more the better.'
    for paragraph in mypres.slides[1].shapes[10].text_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.name = 'Times New Roman'
        
    mypres.slides[1].shapes[10].text_frame.paragraphs[3].text = 'If you START investing and stick to the plan of monthly SIP of '+str(format_currency(sip/12, 'INR', locale='en_IN'))+ ' and earn avg. annual return of '+str(aror*100)+'% on equity portion and '+str(round(figi*100,1))+' % on debt portion , this is how your portfolio will grow- '
    for paragraph in mypres.slides[1].shapes[10].text_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.name = 'Times New Roman'    
        
    srn = counter() # generates a srn number

    mypres.slides[0].shapes[12].text_frame.paragraphs[0].text = 'SRN: ' + str(srn)+ '\nTransaction ID: '+str(int(mobile)-1000)+'\nReport Created:  '+str(date)
    for paragraph in mypres.slides[0].shapes[12].text_frame.paragraphs:
        paragraph.font.size = Pt(8)
        paragraph.font.name = 'Times New Roman'
        
    mypres.slides[1].shapes[7].text_frame.paragraphs[0].text = 'SRN: '+ str(srn)+ '\nTransaction ID: '+str(int(mobile)-1000)+'\nReport Created:  '+str(date)
    for paragraph in mypres.slides[1].shapes[7].text_frame.paragraphs:
        paragraph.font.size = Pt(8)
        paragraph.font.name = 'Times New Roman'
        
        
    def delete_paragraph(paragraph):
        p = paragraph._p
        parent_element = p.getparent()
        parent_element.remove(p)
        
    if data_row[14] != 'less than 10%':
        delete_paragraph(mypres.slides[1].shapes[12].text_frame.paragraphs[1])
        paragraph.font.name = 'Times New Roman'    
        

    mypres.save('base.pptx')




# FINAL CONVERSION TO PDF

def to_pdf(folder_path, base_presentation, name, code):
    '''
    Converts the pptx to pdf

    Parameters
    ----------
    folder_path : path, string
        the path where our pptx 'base.pptx' is.
    base_presentation : path to pptx, string
        the 'base.pptx'.
    name : str
        name.
    code : code
        code.

    Returns
    -------
    Saves a pdf file into the directory.

    '''
    

    folder_path += "\\"
    
    input_file_paths = os.listdir(folder_path)
    
    #Convert each file
    for input_file_name in input_file_paths:
    
        if not input_file_name.lower().endswith((".ppt", ".pptx")):
            continue
        
        # Create input file path
        input_file_path = os.path.join(folder_path, "base.pptx")
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        
        powerpoint.Visible = 1
        
        slides = powerpoint.Presentations.Open(input_file_path)
        
        file_name = os.path.splitext(input_file_name)[0]
        
        output_file_path = os.path.join(folder_path, str(name)+str(code) + ".pdf")
        
        # Save as PDF (formatType = 32)
        slides.SaveAs(output_file_path, 32)
        
        # Close the slide deck
        slides.Close()
    powerpoint.Quit()



